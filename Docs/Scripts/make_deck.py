"""
Generate fl.AI Gene Classification slide deck for senior non-technical research scientists.
Run:  python3 make_deck.py
Output: fl.AI_Gene_Classification_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x2B, 0x55)   # deep navy
BLUE    = RGBColor(0x1A, 0x5F, 0xA8)   # mid blue
TEAL    = RGBColor(0x00, 0x91, 0x87)   # teal accent
GREEN   = RGBColor(0x2E, 0x7D, 0x32)   # success green
ORANGE  = RGBColor(0xE6, 0x5C, 0x00)   # warm orange
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF4, 0xF6, 0xF9)   # light grey bg
MGRAY   = RGBColor(0xB0, 0xBE, 0xC5)   # mid grey text
DKGRAY  = RGBColor(0x37, 0x47, 0x4F)   # dark grey text
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)   # highlight

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=DKGRAY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_label_value(slide, label, value, left, top, label_w=2.2, value_w=2.8,
                    row_h=0.38, label_color=NAVY, value_color=DKGRAY, font_size=13):
    add_text_box(slide, label, left, top, label_w, row_h,
                 font_size=font_size, bold=True, color=label_color)
    add_text_box(slide, value, left + label_w, top, value_w, row_h,
                 font_size=font_size, color=value_color)


def header_band(slide, title, subtitle=None):
    """Dark navy top band with title."""
    add_rect(slide, 0, 0, 13.33, 1.45, NAVY)
    add_text_box(slide, title, 0.35, 0.10, 12.5, 0.80,
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.35, 0.82, 12.5, 0.50,
                     font_size=15, color=MGRAY, align=PP_ALIGN.LEFT)


def section_label(slide, text, left, top, width=3.5):
    """Small teal ALL-CAPS section label."""
    add_text_box(slide, text.upper(), left, top, width, 0.28,
                 font_size=9, bold=True, color=TEAL)


def bullet_box(slide, items, left, top, width, height,
               font_size=13, bullet="•", color=DKGRAY, line_spacing=0.36):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def pill(slide, text, left, top, width, height, bg=BLUE, fg=WHITE, font_size=12, bold=True):
    r = add_rect(slide, left, top, width, height, bg)
    r.adjustments[0] = 50000   # rounded corners (only works for rounded rect shape type 5)
    t = add_text_box(slide, text, left, top, width, height,
                     font_size=font_size, bold=bold, color=fg,
                     align=PP_ALIGN.CENTER)
    return r, t


def arrow_right(slide, left, top, length=0.3, color=MGRAY):
    """Simple right-pointing arrow line."""
    from pptx.util import Pt as ptutil
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(left), Inches(top), Inches(left + length), Inches(top)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


# ── Build slides ─────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Full background split: navy left, white right
add_rect(slide, 0, 0, 7.5, 7.5, NAVY)
add_rect(slide, 7.5, 0, 5.83, 7.5, LGRAY)

# Accent stripe
add_rect(slide, 7.5, 0, 0.06, 7.5, TEAL)

# fl.AI logo text area
add_text_box(slide, "fl.AI", 0.6, 1.4, 6.5, 1.5,
             font_size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "Automated Gene Literature\nClassification Pipeline",
             0.6, 3.0, 6.5, 1.5,
             font_size=26, bold=False, color=MGRAY, align=PP_ALIGN.LEFT)
add_text_box(slide, "Drosophila melanogaster  ·  AI-Assisted  ·  Lab-Scale",
             0.6, 4.6, 6.5, 0.5,
             font_size=14, bold=False, color=TEAL, align=PP_ALIGN.LEFT)

# Right panel – context cards
add_text_box(slide, "Presented to", 8.0, 1.3, 5.0, 0.35,
             font_size=11, color=MGRAY)
add_text_box(slide, "Senior Research Scientists", 8.0, 1.65, 5.0, 0.5,
             font_size=18, bold=True, color=NAVY)

add_rect(slide, 7.9, 2.5, 4.7, 0.5, BLUE)
add_text_box(slide, "The Allada Lab  ·  University of Michigan", 7.9, 2.5, 4.7, 0.5,
             font_size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# Key stats
stats = [
    ("Genes processed",    "Hundreds per run"),
    ("Literature sources", "FlyBase · PubMed · Europe PMC"),
    ("AI Model",           "GPT (OpenAI)"),
    ("Output",             "Structured Excel workbook"),
]
for i, (k, v) in enumerate(stats):
    y = 3.4 + i * 0.7
    add_rect(slide, 7.9, y, 4.7, 0.55, WHITE)
    add_text_box(slide, k, 8.05, y + 0.04, 2.0, 0.45,
                 font_size=11, bold=True, color=NAVY)
    add_text_box(slide, v, 10.1, y + 0.04, 2.5, 0.45,
                 font_size=11, color=DKGRAY)

add_text_box(slide, "2026", 0.6, 6.8, 2, 0.4,
             font_size=12, color=MGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – The Challenge
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "The Challenge", "Manually keeping up with fly gene literature is impractical at scale")

# Two-column layout
# Left: the problem
add_rect(slide, 0.3, 1.65, 5.9, 5.4, LGRAY)
section_label(slide, "The Old Way", 0.55, 1.75)
bullets_left = [
    "A typical gene-expression study yields 50–500 candidate genes",
    "Each gene may have dozens of published papers across PubMed, FlyBase, and Europe PMC",
    "Researchers must search each gene manually, read full-text articles, and judge relevance",
    "Extracting function, phenotype, and reagent data is repetitive and error-prone",
    "Results are inconsistent across lab members and hard to audit or reproduce",
]
bullet_box(slide, bullets_left, 0.55, 2.2, 5.4, 4.5, font_size=13)

# Right: the consequence
add_rect(slide, 6.6, 1.65, 6.4, 5.4, LGRAY)
section_label(slide, "Consequences", 6.85, 1.75)

# Pain cards
pains = [
    (ORANGE, "Weeks of manual curation for one gene set"),
    (ORANGE, "Important papers easily missed"),
    (ORANGE, "No systematic record of evidence reviewed"),
    (ORANGE, "Hard to update when new papers publish"),
    (GREEN,  "→ fl.AI was built to solve exactly this"),
]
for i, (col, text) in enumerate(pains):
    y = 2.2 + i * 0.88
    add_rect(slide, 6.85, y, 5.9, 0.72, col)
    add_text_box(slide, text, 7.05, y + 0.08, 5.5, 0.55,
                 font_size=13, bold=(col == GREEN), color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Solution Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "fl.AI at a Glance", "One automated pipeline from gene list to classified evidence report")

# Central tagline
add_rect(slide, 0.3, 1.6, 12.73, 0.65, TEAL)
add_text_box(slide,
    "Give fl.AI a list of fly genes and research keywords → get back a fully cited, AI-classified Excel report.",
    0.4, 1.65, 12.5, 0.55,
    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three pillars
pillars = [
    (BLUE,   "📚  Literature\nAggregation",
     "Automatically collects all relevant publications from FlyBase, PubMed, and Europe PMC for every gene in the list."),
    (TEAL,   "🧠  AI Evidence\nExtraction",
     "GPT reads full-text papers and extracts gene function, phenotype observations, and genetic reagents in plain language."),
    (GREEN,  "🗂  Structured\nClassification",
     "Each gene is assigned to researcher-defined categories (e.g. sleep, circadian, synapse) with a confidence score and evidence rationale."),
]
for i, (col, title, desc) in enumerate(pillars):
    x = 0.35 + i * 4.28
    add_rect(slide, x, 2.55, 4.0, 1.05, col)
    add_text_box(slide, title, x + 0.15, 2.6, 3.7, 1.0,
                 font_size=17, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(slide, x, 3.62, 4.0, 2.85, LGRAY)
    add_text_box(slide, desc, x + 0.15, 3.7, 3.7, 2.6,
                 font_size=13, color=DKGRAY)

# Bottom strip
add_rect(slide, 0.3, 6.6, 12.73, 0.6, NAVY)
add_text_box(slide,
    "Designed to run on the Allada Lab HPC (turbo-server)  ·  Fully resumable  ·  No manual steps required",
    0.4, 6.65, 12.5, 0.5,
    font_size=13, color=MGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Pipeline Walkthrough
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "How It Works", "An 8-step automated pipeline processes each gene from symbol to classified report")

steps = [
    (NAVY,  "1", "Gene Symbol\nConversion",  "Gene symbols → FlyBase\nFBgn identifiers"),
    (BLUE,  "2", "FlyBase\nReferences",      "Pull known publication\nIDs from FlyBase"),
    (BLUE,  "3", "PubMed\nSearch",           "Query PubMed with\ngene name variants"),
    (BLUE,  "4", "Europe PMC\nSearch",       "Expand to international\nopen-access papers"),
    (TEAL,  "5", "Rank &\nDeduplicate",      "Score, sort, and cap\nthe reference list"),
    (TEAL,  "6", "Keyword\nFiltering",       "Retain only papers\nmatching your terms"),
    (GREEN, "7", "Full-Text\nRetrieval",     "Download PDFs / XML\nwhen available"),
    (GREEN, "8", "AI Analysis\n& Export",    "GPT extracts evidence;\nclassify & export Excel"),
]

box_w, box_h = 1.42, 1.6
gap = 0.1
start_x = 0.25

for i, (col, num, title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    # number circle background
    add_rect(slide, x, 1.55, box_w, 0.42, col)
    add_text_box(slide, f"Step {num}", x, 1.55, box_w, 0.42,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 1.97, box_w, box_h - 0.42, LGRAY)
    add_text_box(slide, title, x + 0.05, 2.02, box_w - 0.1, 0.65,
                 font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x + 0.05, 2.68, box_w - 0.1, 0.85,
                 font_size=10, color=DKGRAY, align=PP_ALIGN.CENTER)

# Arrow between steps using thin lines
for i in range(len(steps) - 1):
    ax = start_x + (i + 1) * (box_w + gap) - gap + 0.02
    ay = 1.55 + (box_h / 2) - 0.02
    add_text_box(slide, "▶", ax - 0.05, ay - 0.04, 0.18, 0.22,
                 font_size=8, color=MGRAY, align=PP_ALIGN.CENTER)

# Example numbers at bottom
add_rect(slide, 0.25, 3.85, 12.83, 0.4, NAVY)
add_text_box(slide, "Example run  (1 gene, keywords: sleep · circadian · synapse)",
             0.35, 3.88, 12.5, 0.35, font_size=12, bold=True, color=WHITE)

ex_data = [
    ("FlyBase refs", "9"),
    ("PubMed refs",  "14"),
    ("EuropePMC",    "11"),
    ("After dedup",  "21"),
    ("After limit",  "15"),
    ("Keyword pass", "10"),
    ("Full-text got","7"),
    ("Summaries",    "6"),
]
for i, (lbl, val) in enumerate(ex_data):
    x = 0.35 + i * 1.58
    add_rect(slide, x, 4.35, 1.42, 0.85, LGRAY)
    add_text_box(slide, val, x, 4.38, 1.42, 0.48,
                 font_size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, lbl, x, 4.82, 1.42, 0.35,
                 font_size=9, color=DKGRAY, align=PP_ALIGN.CENTER)

# Footer note
add_text_box(slide,
    "Per-gene processing typically takes 30–60 seconds  ·  Full gene sets run as overnight batches  ·  Automatically resumes if interrupted",
    0.25, 5.35, 12.83, 0.4,
    font_size=10, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Literature Coverage
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "Comprehensive Literature Coverage",
            "fl.AI searches three complementary databases to maximise recall")

sources = [
    (NAVY,  "FlyBase",
     "The authoritative database for Drosophila genetics. Contains curated gene-publication links validated by FlyBase curators.",
     ["Pre-computed gene → paper mappings", "Covers experimental data back to the 1990s", "Used as the gold-standard starting point"]),
    (BLUE,  "PubMed / NCBI",
     "The world's largest biomedical literature database. fl.AI queries PubMed with multiple gene name variants and synonyms.",
     ["Catches papers not yet in FlyBase", "Includes preprints and recent publications", "Uses gene symbol variants & Greek-letter forms"]),
    (TEAL,  "Europe PMC",
     "European counterpart to PubMed with strong open-access full-text coverage, enabling higher rates of full-text retrieval.",
     ["Extends open-access full-text availability", "International journal coverage", "Provides structured XML for clean text extraction"]),
]

for i, (col, name, desc, bullets) in enumerate(sources):
    x = 0.3 + i * 4.28
    add_rect(slide, x, 1.6, 4.1, 0.65, col)
    add_text_box(slide, name, x + 0.15, 1.65, 3.8, 0.55,
                 font_size=20, bold=True, color=WHITE)
    add_rect(slide, x, 2.25, 4.1, 1.5, LGRAY)
    add_text_box(slide, desc, x + 0.15, 2.32, 3.8, 1.35,
                 font_size=12, color=DKGRAY)
    add_rect(slide, x, 3.77, 4.1, 2.5, WHITE)
    for j, b in enumerate(bullets):
        add_text_box(slide, f"✓  {b}", x + 0.15, 3.82 + j * 0.72, 3.8, 0.65,
                     font_size=12, color=GREEN)

# Merging note
add_rect(slide, 0.3, 6.4, 12.73, 0.75, NAVY)
add_text_box(slide,
    "After collection, fl.AI deduplicates across all three sources and ranks references by how many sources agree on them — "
    "papers cited by all three sources rank highest.",
    0.5, 6.45, 12.3, 0.65,
    font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – AI Evidence Extraction
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "AI-Powered Evidence Extraction",
            "GPT reads each paper and produces structured summaries — not just keyword matches")

# Left: what GPT does
add_rect(slide, 0.3, 1.6, 5.9, 5.55, LGRAY)
section_label(slide, "What the AI extracts", 0.55, 1.72)

extractions = [
    (BLUE, "Gene Function",
     "What biological role does this gene play? What molecular mechanisms are described?"),
    (TEAL, "Phenotype Evidence",
     "What happens to the fly when this gene is mutated, overexpressed, or knocked down?"),
    (GREEN, "Genetic Reagents",
     "Which fly stocks, RNAi lines, or transgenic tools were used to study this gene?"),
]
for i, (col, label, desc) in enumerate(extractions):
    y = 2.15 + i * 1.6
    add_rect(slide, 0.5, y, 0.15, 1.35, col)
    add_text_box(slide, label, 0.8, y + 0.05, 5.1, 0.4,
                 font_size=14, bold=True, color=col)
    add_text_box(slide, desc, 0.8, y + 0.45, 5.1, 0.85,
                 font_size=12, color=DKGRAY)

# Right: example output box
add_rect(slide, 6.5, 1.6, 6.55, 5.55, LGRAY)
section_label(slide, "Example extracted summary", 6.75, 1.72)

add_rect(slide, 6.7, 2.1, 6.15, 0.38, NAVY)
add_text_box(slide, "Gene: period (per)  ·  Paper: PMCID 3456789", 6.75, 2.12, 6.0, 0.33,
             font_size=11, bold=True, color=WHITE)

example_text = (
    "FUNCTION:\n"
    "period encodes a core component of the circadian clock. The PER protein dimerises "
    "with TIM and undergoes rhythmic nuclear translocation, repressing CLK/CYC-driven "
    "transcription of clock genes in a ~24 h feedback loop.\n\n"
    "PHENOTYPE:\n"
    "Loss-of-function per01 flies are behaviourally arrhythmic. Overexpression shortens "
    "or lengthens free-running period length depending on dosage. Sleep is severely "
    "disrupted in null mutants.\n\n"
    "REAGENTS IDENTIFIED:\n"
    "• per01 null allele (FlyBase FBal0000851)\n"
    "• UAS-per (Bloomington #7127)\n"
    "• tim-GAL4 driver line"
)

add_text_box(slide, example_text, 6.7, 2.55, 6.15, 4.4,
             font_size=10.5, color=DKGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Classification Engine
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "Gene Classification",
            "Every gene receives a category assignment, confidence score, and evidence-backed rationale")

# How classification works
add_rect(slide, 0.3, 1.6, 12.73, 0.55, LGRAY)
add_text_box(slide,
    "Researchers define their own categories (e.g.  sleep  ·  circadian  ·  synapse  ·  metabolism).  "
    "fl.AI assigns each gene to zero, one, or multiple categories.",
    0.5, 1.64, 12.3, 0.47, font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

# Classification output table mockup
headers = ["Gene Symbol", "FlyBase ID", "Category", "Confidence", "Rationale (excerpt)"]
col_widths = [1.6, 1.5, 1.6, 1.3, 6.5]
rows = [
    ["period", "FBgn0003068", "sleep · circadian", "94",
     "Multiple papers document severe sleep disruption and arrhythmia in per mutants…"],
    ["timeless", "FBgn0014396", "circadian", "91",
     "TIM is a canonical circadian clock component; no direct sleep evidence found…"],
    ["redeye", "FBgn0031722", "sleep", "78",
     "Loss-of-function increases total sleep duration; no circadian phenotype reported…"],
    ["bruchpilot", "FBgn0025682", "synapse", "88",
     "Essential active zone protein; null mutants show impaired synaptic vesicle release…"],
    ["sifamide", "FBgn0039170", "sleep · metabolism", "72",
     "Neuropeptide modulating both sleep need and feeding; moderate evidence for both…"],
]

# Header row
add_rect(slide, 0.3, 2.35, 12.73, 0.42, NAVY)
x_pos = 0.35
for j, (h, w) in enumerate(zip(headers, col_widths)):
    add_text_box(slide, h, x_pos, 2.38, w - 0.05, 0.35,
                 font_size=11, bold=True, color=WHITE)
    x_pos += w

for i, row in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    add_rect(slide, 0.3, 2.77 + i * 0.65, 12.73, 0.62, bg)
    x_pos = 0.35
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        cell_color = BLUE if j == 2 else DKGRAY
        bold = j == 2
        add_text_box(slide, cell, x_pos, 2.8 + i * 0.65, w - 0.05, 0.55,
                     font_size=10.5, bold=bold, color=cell_color)
        x_pos += w

# Confidence explanation
add_rect(slide, 0.3, 6.15, 12.73, 0.95, LGRAY)
section_label(slide, "Confidence Score", 0.55, 6.2)
add_text_box(slide,
    "0–100 score reflecting the strength and consistency of published evidence.  "
    "A score ≥ 80 indicates strong, multi-paper support.  "
    "Scores of 50–79 reflect partial or indirect evidence.  "
    "Below 50: weak or ambiguous support — treat as a lead, not a conclusion.",
    0.55, 6.42, 12.3, 0.6,
    font_size=12, color=DKGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Output Workbook
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "The Output: Excel Workbook",
            "One file, four structured sheets — ready to filter, sort, and share")

sheets = [
    (NAVY,  "Sheet 1",  "Gene Set",
     "The original input gene list with all metadata preserved. Serves as the master index for cross-referencing."),
    (BLUE,  "Sheet 2",  "Classification",
     "The primary deliverable. Each gene with its assigned categories, confidence score, rationale, and the full evidence block used for the decision."),
    (TEAL,  "Sheet 3",  "Reference Summaries",
     "Every paper reviewed, with title, year, journal, author, abstract, full evidence summary, and source database."),
    (GREEN, "Sheet 4",  "Reagents",
     "Genetic tools discovered in the literature: fly stocks, RNAi lines, transgenic constructs. Includes stock IDs and evidence snippets."),
]

for i, (col, badge, name, desc) in enumerate(sheets):
    y = 1.65 + i * 1.35
    add_rect(slide, 0.3, y, 1.0, 1.18, col)
    add_text_box(slide, badge, 0.3, y + 0.05, 1.0, 0.38,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, name, 0.3, y + 0.45, 1.0, 0.65,
                 font_size=9, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, 1.32, y, 11.7, 1.18, LGRAY)
    add_text_box(slide, name, 1.5, y + 0.04, 3.0, 0.42,
                 font_size=15, bold=True, color=col)
    add_text_box(slide, desc, 1.5, y + 0.46, 11.2, 0.65,
                 font_size=12.5, color=DKGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Reagent Discovery
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "Bonus: Automatic Reagent Discovery",
            "fl.AI also catalogs the genetic tools used in every reviewed paper")

# Intro
add_text_box(slide,
    "For each gene, fl.AI identifies every fly stock, RNAi line, or transgenic construct mentioned "
    "in the reviewed papers — saving hours of reagent-hunting for follow-up experiments.",
    0.35, 1.65, 12.6, 0.65,
    font_size=14, color=DKGRAY)

# Reagent table
r_headers = ["Gene", "Stock ID", "Collection", "Reagent Type", "Functional Validity", "Evidence Snippet"]
r_widths   = [1.3, 1.4, 1.5, 1.5, 1.6, 5.7]
r_rows = [
    ["period", "FBal0000851", "FlyBase",      "Null allele",    "Validated",   "per01 flies show complete behavioural arrhythmia in LD and DD conditions."],
    ["period", "#7127",       "Bloomington",  "UAS transgene",  "Validated",   "UAS-per overexpression rescues arrhythmia in per01 background."],
    ["timeless", "FBal0015611","FlyBase",     "Null allele",    "Validated",   "tim01 null: arrhythmic locomotor activity; TIM protein absent."],
    ["redeye",  "BL#23700",   "Bloomington",  "RNAi line",      "Partial",     "Knockdown in LNv neurons increases sleep duration by ~40 min."],
]

add_rect(slide, 0.3, 2.45, 12.73, 0.42, NAVY)
x_pos = 0.35
for h, w in zip(r_headers, r_widths):
    add_text_box(slide, h, x_pos, 2.48, w - 0.05, 0.35,
                 font_size=11, bold=True, color=WHITE)
    x_pos += w

for i, row in enumerate(r_rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    add_rect(slide, 0.3, 2.87 + i * 0.6, 12.73, 0.57, bg)
    x_pos = 0.35
    for j, (cell, w) in enumerate(zip(row, r_widths)):
        add_text_box(slide, cell, x_pos, 2.9 + i * 0.6, w - 0.05, 0.5,
                     font_size=10, color=DKGRAY)
        x_pos += w

# Value proposition
add_rect(slide, 0.3, 5.35, 12.73, 1.7, LGRAY)
section_label(slide, "Why this matters", 0.55, 5.45)
benefits = [
    "Instantly know which validated stocks exist for your candidate genes",
    "Distinguish established null alleles from newer conditional tools",
    "Identify reagents available in Bloomington or VDRC stock collections",
    "Accelerate experimental design by linking evidence directly to tools",
]
for i, b in enumerate(benefits):
    x = 0.5 + (i % 2) * 6.3
    y = 5.82 + (i // 2) * 0.55
    add_text_box(slide, f"✓  {b}", x, y, 5.9, 0.5, font_size=12, color=DKGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Technical Infrastructure
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "Infrastructure & Reliability",
            "Built for unattended, large-scale overnight runs on the lab HPC")

items = [
    (NAVY,  "HPC Integration",
     "Runs on turbo-server (umms-rallada NFS mount). Shared caches avoid redundant API calls across different gene sets run by any lab member."),
    (BLUE,  "Batch Checkpointing",
     "Genes are processed in batches of 50. Progress is saved after every batch. If the job is interrupted (network drop, timeout), it resumes from the last completed batch — no work is lost."),
    (TEAL,  "Smart Caching",
     "PubMed metadata and full-text retrieval methods are cached in shared TSV files. Rerunning a gene set with updated keywords reuses existing downloads, dramatically reducing API usage."),
    (GREEN, "Full-Text Retrieval Strategy",
     "fl.AI tries four methods in sequence — PMC OA PDF, PMC XML, Europe PMC API, and Unpaywall — before falling back to title + abstract only. This maximises the evidence base without paywalled content."),
    (ORANGE,"Reproducibility",
     "Every classification run records which papers were used, which model produced the output, and what keywords drove the analysis. The Excel report is a complete audit trail."),
]

for i, (col, title, desc) in enumerate(items):
    row = i // 2
    col_idx = i % 2
    if i == 4:   # last item spans full width
        x, w = 0.3, 12.73
    else:
        x = 0.3 + col_idx * 6.5
        w = 6.2
    y = 1.65 + row * 1.65
    add_rect(slide, x, y, 0.12, 1.35, col)
    add_text_box(slide, title, x + 0.25, y + 0.05, w - 0.35, 0.38,
                 font_size=14, bold=True, color=col)
    add_text_box(slide, desc, x + 0.25, y + 0.45, w - 0.35, 0.82,
                 font_size=12, color=DKGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Summary & Impact
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
header_band(slide, "Summary: What fl.AI Delivers",
            "From gene list to evidence-based classification — fully automated")

# Before / After
add_rect(slide, 0.3, 1.6, 5.9, 5.4, LGRAY)
add_rect(slide, 0.3, 1.6, 5.9, 0.5, ORANGE)
add_text_box(slide, "Without fl.AI", 0.45, 1.65, 5.6, 0.4,
             font_size=16, bold=True, color=WHITE)
before = [
    "Days to weeks of manual literature searching",
    "Inconsistent evidence standards across lab members",
    "Easy to miss important or recent papers",
    "No structured record of what was reviewed",
    "Difficult to update when new papers are published",
    "Reagent discovery requires separate effort",
]
bullet_box(slide, before, 0.45, 2.2, 5.55, 4.5, font_size=12.5,
           bullet="✗", color=RGBColor(0xC6, 0x28, 0x28))

add_rect(slide, 6.6, 1.6, 6.43, 5.4, LGRAY)
add_rect(slide, 6.6, 1.6, 6.43, 0.5, GREEN)
add_text_box(slide, "With fl.AI", 6.75, 1.65, 6.1, 0.4,
             font_size=16, bold=True, color=WHITE)
after = [
    "Overnight run processes hundreds of genes",
    "Consistent AI-driven evidence extraction for every gene",
    "Multi-source search with priority ranking",
    "Full Excel audit trail with cited references",
    "Rerun with updated keywords in minutes",
    "Reagent sheet auto-populated from reviewed papers",
]
bullet_box(slide, after, 6.75, 2.2, 6.1, 4.5, font_size=12.5,
           bullet="✓", color=GREEN)

# Arrow between
add_text_box(slide, "▶▶", 6.22, 4.15, 0.4, 0.5,
             font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Thank You / Questions
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
add_rect(slide, 0, 0, 13.33, 0.18, TEAL)
add_rect(slide, 0, 7.32, 13.33, 0.18, TEAL)

add_text_box(slide, "Thank You", 0.5, 1.5, 12.33, 1.3,
             font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Questions & Discussion", 0.5, 2.85, 12.33, 0.7,
             font_size=24, color=MGRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 3.0, 3.8, 7.33, 0.06, TEAL)

add_text_box(slide, "fl.AI is open for use by any Allada Lab member", 0.5, 4.1, 12.33, 0.5,
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide,
    "Repository:  fl.AI-gene-classification  ·  Runs on turbo-server\n"
    "Contact: aadish98@gmail.com for access, questions, or feature requests",
    0.5, 4.65, 12.33, 0.85,
    font_size=14, color=MGRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, "2026  ·  Allada Lab  ·  University of Michigan",
             0.5, 6.9, 12.33, 0.4,
             font_size=11, color=MGRAY, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/Users/aadishms/Desktop/Projects.nosync/fl.AI-gene-classification/fl.AI_Gene_Classification_Deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
