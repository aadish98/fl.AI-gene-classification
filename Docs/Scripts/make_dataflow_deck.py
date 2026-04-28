"""
Generate the fl.AI Data Flow Deck — a non-technical, code-free walkthrough
of how data moves through the gene-classification pipeline.

Output: fl.AI_Data_Flow_Deck.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Palette ──────────────────────────────────────────────────────────────────
INK        = RGBColor(0x0E, 0x1F, 0x33)   # deep ink navy (dominant dark)
INK_SOFT   = RGBColor(0x1A, 0x2D, 0x46)   # softer ink
CREAM      = RGBColor(0xF6, 0xF1, 0xE6)   # warm paper cream (dominant light)
CREAM_DEEP = RGBColor(0xEC, 0xE3, 0xCE)   # deeper cream for striping
TEAL       = RGBColor(0x0E, 0x83, 0x88)   # process / transformation accent
TEAL_DEEP  = RGBColor(0x0A, 0x65, 0x6A)   # deeper teal
CORAL      = RGBColor(0xF7, 0x6F, 0x53)   # sharp data accent
AMBER      = RGBColor(0xE0, 0xB3, 0x41)   # secondary highlight
SLATE      = RGBColor(0x3B, 0x44, 0x52)   # body text on cream
SLATE_LITE = RGBColor(0x6B, 0x73, 0x82)   # secondary muted text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK_MUTED  = RGBColor(0xC8, 0xCE, 0xD8)   # muted text on ink

HEADER_FONT = "Georgia"
BODY_FONT   = "Calibri"


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.adjustments[0] = radius
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_text(slide, text, x, y, w, h, *, size=14, bold=False, italic=False,
             color=SLATE, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             font=BODY_FONT, char_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if char_spacing is not None:
        from pptx.oxml.ns import qn
        rpr = r._r.get_or_add_rPr()
        rpr.set("spc", str(char_spacing))
    return tb


def add_lines(slide, lines, x, y, w, h, *, size=13, color=SLATE,
              align=PP_ALIGN.LEFT, font=BODY_FONT, line_gap=4, bold=False,
              italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(line_gap)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_bullets(slide, items, x, y, w, h, *, size=13, color=SLATE,
                marker="—", marker_color=None, font=BODY_FONT, line_gap=6):
    """Custom dash markers for elegance (no unicode bullet)."""
    if marker_color is None:
        marker_color = TEAL
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_gap)
        rm = p.add_run()
        rm.text = f"{marker}  "
        rm.font.size = Pt(size)
        rm.font.bold = True
        rm.font.name = font
        rm.font.color.rgb = marker_color
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(size)
        rt.font.name = font
        rt.font.color.rgb = color
    return tb


def eyebrow(slide, text, x, y, w=4.0, color=TEAL):
    """ALL-CAPS small label."""
    add_text(slide, text.upper(), x, y, w, 0.28,
             size=10, bold=True, color=color, font=BODY_FONT,
             char_spacing=300)


def page_number(slide, n, total):
    add_text(slide, f"{n:02d} / {total:02d}", 12.55, 7.05, 0.7, 0.3,
             size=9, color=SLATE_LITE, align=PP_ALIGN.RIGHT, font=BODY_FONT)


def footer_brand(slide, color=SLATE_LITE):
    add_text(slide, "fl.AI  ·  Data Flow Documentation",
             0.6, 7.05, 6.0, 0.3,
             size=9, color=color, font=BODY_FONT, char_spacing=120)


def horizontal_rule(slide, x, y, w, color=CREAM_DEEP, thick=Pt(0.75)):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                      Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = thick


def chip(slide, text, x, y, w, h, *, fill=CREAM_DEEP, fg=INK,
         size=11, bold=True, font=BODY_FONT, radius=0.5):
    s = add_round_rect(slide, x, y, w, h, fill, radius=radius)
    add_text(slide, text, x, y, w, h,
             size=size, bold=bold, color=fg, font=font,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return s


def arrow_glyph(slide, x, y, w=0.45, h=0.45, color=CORAL, size=22):
    """Decorative arrow glyph between flow stages."""
    add_text(slide, "→", x, y, w, h,
             size=size, bold=True, color=color, font=BODY_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def stage_header(slide, stage_num, title, kicker, total=7):
    """Standard header for stage slides 3–9."""
    # Cream background already set on slide
    # Stage badge: small ink chip top-left
    add_round_rect(slide, 0.6, 0.55, 1.6, 0.42, INK, radius=0.5)
    add_text(slide, f"Stage {stage_num}  /  {total}",
             0.6, 0.55, 1.6, 0.42,
             size=11, bold=True, color=CREAM, font=BODY_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
             char_spacing=150)
    # Title
    add_text(slide, title, 0.6, 1.08, 12.1, 0.85,
             size=36, bold=True, color=INK, font=HEADER_FONT)
    # Kicker / one-line subtitle
    add_text(slide, kicker, 0.62, 1.95, 12.1, 0.45,
             size=15, italic=True, color=SLATE_LITE, font=HEADER_FONT)
    # Thin horizontal rule under header
    horizontal_rule(slide, 0.6, 2.55, 12.1, color=CREAM_DEEP, thick=Pt(1.0))


def transformation_strip(slide, y, *, before, label, after,
                         before_fill=CREAM_DEEP, label_fill=TEAL,
                         after_fill=CORAL, after_fg=WHITE):
    """A horizontal INPUT → [PROCESS] → OUTPUT strip."""
    chip(slide, before, 0.6, y, 4.0, 0.55,
         fill=before_fill, fg=INK, size=12, radius=0.6)
    arrow_glyph(slide, 4.65, y, w=0.5, h=0.55, color=SLATE_LITE, size=22)
    chip(slide, label, 5.2, y, 3.0, 0.55,
         fill=label_fill, fg=WHITE, size=12, radius=0.6)
    arrow_glyph(slide, 8.25, y, w=0.5, h=0.55, color=SLATE_LITE, size=22)
    chip(slide, after, 8.8, y, 3.9, 0.55,
         fill=after_fill, fg=after_fg, size=12, radius=0.6)


# ── Flowchart helpers (organism-specific slides) ────────────────────────────

def flowchart_card(slide, x, y, w, h, *, num, title, desc,
                   accent=TEAL, body_size=9, title_size=12):
    """Single flowchart node: white card, accent stripe, number, title, desc."""
    add_round_rect(slide, x, y, w, h, WHITE, radius=0.04)
    add_rect(slide, x, y, 0.12, h, accent)
    # number, top-right corner
    add_text(slide, f"{num:02d}", x + w - 0.55, y + 0.10, 0.45, 0.30,
             size=11, bold=True, color=accent, font=HEADER_FONT,
             align=PP_ALIGN.RIGHT)
    # title
    add_text(slide, title, x + 0.27, y + 0.10, w - 0.85, 0.32,
             size=title_size, bold=True, color=INK, font=HEADER_FONT)
    # description body
    add_text(slide, desc, x + 0.27, y + 0.46, w - 0.42, h - 0.55,
             size=body_size, color=SLATE, font=BODY_FONT)


def snake_flowchart(slide, *, steps, cols, accent,
                    grid_x=0.6, grid_y=2.85, grid_w=12.1,
                    box_h=1.30, gap_x=0.20, row_gap=0.55):
    """
    Render a snake-pattern flowchart (row 0 left→right, row 1 right→left, etc.)
    with arrows between consecutive cards.
    """
    n = len(steps)
    box_w = (grid_w - (cols - 1) * gap_x) / cols

    # Compute (x, y) positions per step
    positions = []
    for i in range(n):
        row = i // cols
        col_in_row = i % cols
        if row % 2 == 0:
            col = col_in_row
        else:
            col = (cols - 1) - col_in_row
        x = grid_x + col * (box_w + gap_x)
        y = grid_y + row * (box_h + row_gap)
        positions.append((x, y, row, col))

    # Draw cards
    for (x, y, _row, _col), (num, title, desc) in zip(positions, steps):
        flowchart_card(slide, x, y, box_w, box_h,
                       num=num, title=title, desc=desc, accent=accent)

    # Draw connectors between consecutive cards
    for i in range(n - 1):
        x1, y1, r1, c1 = positions[i]
        x2, y2, r2, c2 = positions[i + 1]
        if r1 == r2:
            # horizontal arrow between cards in the same row
            arrow_y = y1 + box_h / 2 - 0.18
            if c2 > c1:
                # left → right
                ax = x1 + box_w
                add_text(slide, "→", ax, arrow_y, gap_x, 0.36,
                         size=18, bold=True, color=accent, font=BODY_FONT,
                         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            else:
                # right → left
                ax = x2 + box_w
                add_text(slide, "←", ax, arrow_y, gap_x, 0.36,
                         size=18, bold=True, color=accent, font=BODY_FONT,
                         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        else:
            # down arrow (same column), bridging the row gap
            ax = x1 + box_w / 2 - 0.18
            ay = y1 + box_h
            add_text(slide, "↓", ax, ay, 0.36, row_gap,
                     size=18, bold=True, color=accent, font=BODY_FONT,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def organism_flowchart_slide(prs, blank, *, slide_num, total,
                             organism_label, title, subtitle,
                             accent, steps, cols,
                             input_label, output_label):
    """Build one organism-specific flowchart slide (fly / human / mouse)."""
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, CREAM)

    # Header — eyebrow + title + organism badge on the right
    eyebrow(slide, f"Flowchart  ·  {organism_label} mode",
            0.6, 0.6, w=8.0, color=accent)
    add_text(slide, title, 0.6, 0.95, 9.5, 0.85,
             size=32, bold=True, color=INK, font=HEADER_FONT)
    add_text(slide, subtitle, 0.6, 1.85, 12.1, 0.55,
             size=14, italic=True, color=SLATE_LITE, font=HEADER_FONT)

    # Organism badge, right-aligned in header
    badge_w = 1.8
    badge_x = 13.33 - 0.6 - badge_w
    add_round_rect(slide, badge_x, 0.95, badge_w, 0.55, accent, radius=0.5)
    add_text(slide, organism_label.upper(),
             badge_x, 0.95, badge_w, 0.55,
             size=14, bold=True, color=WHITE, font=HEADER_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
             char_spacing=200)

    # Thin rule under header
    horizontal_rule(slide, 0.6, 2.55, 12.1, color=CREAM_DEEP, thick=Pt(1.0))

    # Snake flowchart body
    snake_flowchart(slide, steps=steps, cols=cols, accent=accent)

    # Bottom output band — input → output summary
    band_y = 6.20
    add_round_rect(slide, 0.6, band_y, 12.1, 0.85, INK, radius=0.04)
    # eyebrow
    add_text(slide, "INPUT", 0.85, band_y + 0.10, 1.2, 0.28,
             size=10, bold=True, color=AMBER, font=BODY_FONT,
             char_spacing=200)
    add_text(slide, input_label, 0.85, band_y + 0.40, 5.6, 0.40,
             size=12, color=CREAM, font=HEADER_FONT,
             valign=MSO_ANCHOR.TOP)
    # divider
    add_text(slide, "→", 6.5, band_y + 0.20, 0.5, 0.5,
             size=20, bold=True, color=accent, font=BODY_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "OUTPUT", 7.1, band_y + 0.10, 1.2, 0.28,
             size=10, bold=True, color=AMBER, font=BODY_FONT,
             char_spacing=200)
    add_text(slide, output_label, 7.1, band_y + 0.40, 5.6, 0.40,
             size=12, color=CREAM, font=HEADER_FONT,
             valign=MSO_ANCHOR.TOP)

    footer_brand(slide)
    page_number(slide, slide_num, total)


def stat_card(slide, x, y, w, h, *, value, label, note=None, accent=TEAL):
    """Compact metric card used by the real-data walkthrough slides."""
    add_round_rect(slide, x, y, w, h, WHITE, radius=0.04)
    add_rect(slide, x, y, w, 0.12, accent)
    add_text(slide, value, x + 0.18, y + 0.22, w - 0.36, 0.45,
             size=25, bold=True, color=accent, font=HEADER_FONT)
    add_text(slide, label, x + 0.18, y + 0.72, w - 0.36, 0.32,
             size=11, bold=True, color=INK, font=BODY_FONT,
             char_spacing=80)
    if note:
        add_text(slide, note, x + 0.18, y + 1.08, w - 0.36, h - 1.16,
                 size=10, color=SLATE, font=BODY_FONT)


def add_simple_table(slide, *, x, y, col_ws, row_h, headers, rows,
                     accent=TEAL, font_size=9.2, header_size=9.5):
    """Small visual table for example result rows."""
    table_w = sum(col_ws)
    add_round_rect(slide, x, y, table_w, row_h * (len(rows) + 1),
                   WHITE, radius=0.03)
    add_rect(slide, x, y, table_w, row_h, INK)
    cx = x
    for header, w in zip(headers, col_ws):
        add_text(slide, header, cx + 0.08, y + 0.06, w - 0.16, row_h - 0.08,
                 size=header_size, bold=True, color=CREAM, font=BODY_FONT,
                 valign=MSO_ANCHOR.MIDDLE)
        cx += w

    for r_idx, row in enumerate(rows):
        ry = y + row_h * (r_idx + 1)
        if r_idx % 2 == 0:
            add_rect(slide, x, ry, table_w, row_h, CREAM)
        cx = x
        for c_idx, (val, w) in enumerate(zip(row, col_ws)):
            color = INK
            bold = c_idx in (0, 2, 3)
            add_text(slide, str(val), cx + 0.08, ry + 0.06, w - 0.16, row_h - 0.08,
                     size=font_size, bold=bold, color=color, font=BODY_FONT,
                     valign=MSO_ANCHOR.MIDDLE)
            cx += w


def example_slide_header(slide, eyebrow_text, title, subtitle, *,
                         slide_num, total, accent=TEAL):
    """Shared header for Tx-Omics_Conserved example slides."""
    add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
    eyebrow(slide, eyebrow_text, 0.6, 0.6, w=8.5, color=accent)
    add_text(slide, title, 0.6, 0.95, 12.1, 0.72,
             size=32, bold=True, color=INK, font=HEADER_FONT)
    add_text(slide, subtitle, 0.6, 1.75, 12.1, 0.52,
             size=14, italic=True, color=SLATE_LITE, font=HEADER_FONT)
    horizontal_rule(slide, 0.6, 2.45, 12.1, color=CREAM_DEEP, thick=Pt(1.0))
    footer_brand(slide)
    page_number(slide, slide_num, total)


EXAMPLE_INPUT_GENES = [
    ("AstA-R2", "FBgn0039595"),
    ("BomBc2", "FBgn0034331"),
    ("SIFa", "FBgn0053527"),
    ("Trhn", "FBgn0035187"),
    ("na", "FBgn0002917"),
    ("rumpel", "FBgn0029950"),
    ("unc79", "FBgn0038693"),
    ("unc80", "FBgn0039536"),
]

HUMAN_RESULT_ROWS = [
    ("NALCN", "na", "sleep; circadian", "92"),
    ("UNC80", "unc80", "sleep; circadian", "77"),
    ("UNC79", "unc79", "sleep; circadian", "68"),
    ("TH", "Trhn", "sleep; circadian", "63"),
    ("GALR3", "AstA-R2", "circadian", "72"),
    ("GALR1", "AstA-R2", "circadian", "64"),
    ("GALR2", "AstA-R2", "circadian", "58"),
]

MOUSE_RESULT_ROWS = [
    ("Nalcn", "na", "sleep; circadian", "96"),
    ("Tph2", "Trhn", "sleep; circadian", "95"),
    ("Kiss1r", "AstA-R2", "circadian", "91"),
    ("Galr1", "AstA-R2", "sleep", "90"),
    ("Slc5a6", "rumpel", "sleep; circadian", "86"),
    ("Th", "Trhn", "sleep; circadian", "84"),
    ("Slc5a11", "rumpel", "circadian", "80"),
    ("Tph1", "Trhn", "circadian", "72"),
]


# ── Build presentation ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

TOTAL_STAGES = 7
TOTAL_SLIDES = 17


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, INK)

# Cream panel (decorative motif: data entering the system)
add_rect(slide, 0, 0, 5.6, 7.5, CREAM)
# A vertical coral stripe acting as the "data inlet"
add_rect(slide, 5.6, 0, 0.07, 7.5, CORAL)

# Left panel — eyebrow + title + subtitle (stacked, generous breathing room)
eyebrow(slide, "fl.AI  ·  Pipeline Documentation", 0.7, 1.0, w=5.2, color=TEAL_DEEP)

# Title — two intentional lines, each on its own row at large display size
add_text(slide, "Inside", 0.7, 1.50, 5.0, 1.10,
         size=72, bold=True, color=INK, font=HEADER_FONT)
add_text(slide, "fl.AI",  0.7, 2.55, 5.0, 1.10,
         size=72, bold=True, color=CORAL, font=HEADER_FONT)

# Subtitle, well below the title
add_lines(slide, ["How data moves", "through the pipeline."],
          0.7, 3.85, 5.0, 1.30,
          size=24, italic=True, color=SLATE, font=HEADER_FONT, line_gap=2)

# Small data-flow ribbon as visual motif (well below the subtitle)
ribbon_y = 5.55
ribbon_x = 0.7
chips_data = [
    ("INPUT",     CREAM_DEEP, INK),
    ("RESOLVE",   TEAL,       WHITE),
    ("GATHER",    TEAL,       WHITE),
    ("FILTER",    TEAL,       WHITE),
    ("READ",      TEAL,       WHITE),
    ("CLASSIFY",  TEAL,       WHITE),
    ("DELIVER",   CORAL,      WHITE),
]
gap = 0.04
chip_w = 0.62
total_w = len(chips_data) * chip_w + (len(chips_data) - 1) * gap
# Render as a line of small dots representing the pipeline
for i, (lbl, fill, fg) in enumerate(chips_data):
    cx = ribbon_x + i * (chip_w + gap)
    add_oval(slide, cx, ribbon_y, chip_w, chip_w, fill)
    add_text(slide, lbl, cx - 0.1, ribbon_y + chip_w + 0.05,
             chip_w + 0.2, 0.22,
             size=7, bold=True, color=SLATE, font=BODY_FONT,
             align=PP_ALIGN.CENTER, char_spacing=80)
    if i < len(chips_data) - 1:
        # connector line
        nx = cx + chip_w
        line = slide.shapes.add_connector(
            1, Inches(nx), Inches(ribbon_y + chip_w / 2),
            Inches(nx + gap), Inches(ribbon_y + chip_w / 2))
        line.line.color.rgb = SLATE_LITE
        line.line.width = Pt(1.25)

# Bottom meta
add_text(slide, "A non-technical walkthrough  ·  No code, just the journey",
         0.7, 6.85, 5.0, 0.35,
         size=12, italic=True, color=SLATE_LITE, font=HEADER_FONT)

# Right panel (on ink) — context
add_text(slide, "What this deck covers", 5.95, 1.55, 6.8, 0.5,
         size=14, color=AMBER, font=BODY_FONT, char_spacing=200, bold=True)

right_items = [
    "What raw data the pipeline starts with",
    "How a list of fly genes becomes a list of papers",
    "How papers become structured evidence",
    "How evidence becomes a final classification",
    "What you receive at the end",
]
for i, item in enumerate(right_items):
    y = 2.15 + i * 0.55
    # small coral square
    add_rect(slide, 5.95, y + 0.08, 0.16, 0.16, CORAL)
    add_text(slide, item, 6.25, y, 6.5, 0.45,
             size=15, color=WHITE, font=HEADER_FONT)

# Right footer
add_text(slide, "Allada Lab  ·  University of Michigan  ·  2026",
         5.95, 6.85, 6.8, 0.35,
         size=11, color=INK_MUTED, font=BODY_FONT, char_spacing=120)

page_number(slide, 1, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — At a Glance: end-to-end flow
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)

eyebrow(slide, "The journey, end to end", 0.6, 0.6, w=6.0, color=TEAL_DEEP)
add_text(slide, "From a list of genes to a classified evidence report",
         0.6, 0.95, 12.1, 0.85, size=32, bold=True, color=INK, font=HEADER_FONT)
add_text(slide,
         "Seven stages move the data from raw input to the final report. "
         "Each stage adds context, narrows the field, or transforms one kind "
         "of information into another.",
         0.6, 1.85, 12.1, 0.7, size=14, italic=True,
         color=SLATE_LITE, font=HEADER_FONT)

# Seven stages in a horizontal "metro" diagram
stages = [
    ("01", "Inputs",         "Gene lists",         INK),
    ("02", "Resolve",        "Canonical IDs",      INK_SOFT),
    ("03", "Gather",         "Candidate papers",   TEAL_DEEP),
    ("04", "Filter",         "Relevant papers",    TEAL),
    ("05", "Read",           "Evidence summaries", TEAL),
    ("06", "Classify",       "Category verdicts",  CORAL),
    ("07", "Deliver",        "Excel report",       CORAL),
]

# Layout: cards along the top, horizontal track below
n = len(stages)
left_pad = 0.6
right_pad = 0.6
card_w = (13.33 - left_pad - right_pad - (n - 1) * 0.12) / n
card_h = 1.55
card_y = 3.05

for i, (num, name, payload, color) in enumerate(stages):
    cx = left_pad + i * (card_w + 0.12)
    # card body
    add_round_rect(slide, cx, card_y, card_w, card_h, WHITE, radius=0.04)
    # color top accent (rectangle, not rounded so corners look clean)
    add_rect(slide, cx, card_y, card_w, 0.18, color)
    # number
    add_text(slide, num, cx + 0.18, card_y + 0.30, card_w - 0.36, 0.4,
             size=18, bold=True, color=color, font=HEADER_FONT)
    # name
    add_text(slide, name, cx + 0.18, card_y + 0.70, card_w - 0.36, 0.4,
             size=15, bold=True, color=INK, font=HEADER_FONT)
    # payload
    add_text(slide, payload, cx + 0.18, card_y + 1.10, card_w - 0.36, 0.4,
             size=10, color=SLATE_LITE, font=BODY_FONT)

# Track line connecting bottoms of cards
track_y = card_y + card_h + 0.45
line_x1 = left_pad + 0.4
line_x2 = 13.33 - right_pad - 0.4
line = slide.shapes.add_connector(1,
                                  Inches(line_x1), Inches(track_y),
                                  Inches(line_x2), Inches(track_y))
line.line.color.rgb = SLATE_LITE
line.line.width = Pt(1.25)

# Three flow labels under the track
section_labels = [
    ("Identify the genes",           line_x1,                       2.6),
    ("Find & filter the literature", line_x1 + 4.4,                 4.4),
    ("Read & classify",              line_x1 + 9.0,                 3.5),
]
# small coral dots on track
for label, x_label, w_label in section_labels:
    add_oval(slide, x_label - 0.07, track_y - 0.07, 0.14, 0.14, CORAL)
    add_text(slide, label, x_label - 0.5, track_y + 0.18, w_label, 0.35,
             size=12, italic=True, bold=True, color=SLATE, font=HEADER_FONT,
             align=PP_ALIGN.LEFT)

# Bottom band: "What changes at each stage"
band_y = 5.95
add_round_rect(slide, 0.6, band_y, 12.1, 0.95, INK, radius=0.04)
add_text(slide, "WHAT CHANGES AT EACH STAGE",
         0.85, band_y + 0.10, 11.6, 0.3,
         size=10, bold=True, color=AMBER, font=BODY_FONT, char_spacing=200)
add_text(slide,
         "The shape of the data evolves: gene names become canonical IDs; "
         "IDs fan out into thousands of references; references collapse to "
         "a relevant few; relevant papers become structured evidence; "
         "evidence becomes a single decision per gene.",
         0.85, band_y + 0.40, 11.6, 0.55,
         size=13, italic=True, color=CREAM, font=HEADER_FONT)

footer_brand(slide)
page_number(slide, 2, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Stage 1: Inputs
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 1,
             "The starting material",
             "What the pipeline is given before any work begins.",
             total=TOTAL_STAGES)

# LEFT: input description
eyebrow(slide, "What enters the pipeline", 0.6, 2.85, w=6.0)

add_text(slide,
         "A folder of spreadsheets — one per gene set you want to study.",
         0.6, 3.18, 6.4, 0.55,
         size=18, bold=True, color=INK, font=HEADER_FONT)

add_bullets(slide, [
    "Each spreadsheet lists fly gene symbols in a single column.",
    "Different gene sets can come from different experiments.",
    "Researcher also supplies a short list of keywords (the categories of interest).",
    "Optionally, the researcher chooses a target species: fly, human, or mouse.",
], 0.6, 3.95, 6.4, 2.5, size=13, color=SLATE)

# Tiny note
add_text(slide,
         "No special formatting required — the pipeline reads the column you point it to.",
         0.6, 6.5, 6.4, 0.4,
         size=11, italic=True, color=SLATE_LITE, font=HEADER_FONT)

# RIGHT: visual mock of an input gene set
right_x = 7.3
right_w = 5.4
# Card frame
add_round_rect(slide, right_x, 2.85, right_w, 4.0, WHITE, radius=0.03)
# Header band
add_rect(slide, right_x, 2.85, right_w, 0.45, INK)
add_text(slide, "experiment_42.csv", right_x + 0.18, 2.85, right_w - 0.4, 0.45,
         size=12, bold=True, color=CREAM, font=BODY_FONT,
         valign=MSO_ANCHOR.MIDDLE, char_spacing=120)
add_text(slide, "INPUT", right_x + right_w - 0.85, 2.85, 0.7, 0.45,
         size=9, bold=True, color=AMBER, font=BODY_FONT,
         valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT, char_spacing=200)

# Column header
hdr_y = 3.40
add_rect(slide, right_x, hdr_y, right_w, 0.38, CREAM_DEEP)
add_text(slide, "ext_gene", right_x + 0.25, hdr_y, 2.0, 0.38,
         size=11, bold=True, color=SLATE, font=BODY_FONT,
         valign=MSO_ANCHOR.MIDDLE)
add_text(slide, "(other columns…)", right_x + 2.4, hdr_y, 2.8, 0.38,
         size=10, italic=True, color=SLATE_LITE, font=BODY_FONT,
         valign=MSO_ANCHOR.MIDDLE)

# Mock rows
sample_genes = ["per", "tim", "cry", "Clk", "cyc", "vri", "dco", "Pdf", "redeye"]
for i, g in enumerate(sample_genes):
    row_y = hdr_y + 0.42 + i * 0.32
    if i % 2 == 0:
        add_rect(slide, right_x, row_y, right_w, 0.32, CREAM)
    add_text(slide, g, right_x + 0.25, row_y, 2.0, 0.32,
             size=11, color=INK, font=BODY_FONT,
             valign=MSO_ANCHOR.MIDDLE, bold=True)
    add_text(slide, "—", right_x + 2.4, row_y, 2.8, 0.32,
             size=11, color=SLATE_LITE, font=BODY_FONT,
             valign=MSO_ANCHOR.MIDDLE)

# Keywords callout below the file mock — compact so it clears the page number
kw_y = 6.95
add_text(slide, "Keywords:", right_x, kw_y, 1.05, 0.32,
         size=11, bold=True, color=SLATE, font=BODY_FONT,
         valign=MSO_ANCHOR.MIDDLE)
chip_w = 1.05
chip_gap = 0.08
for i, kw in enumerate(["sleep", "circadian", "synapse"]):
    cx = right_x + 1.05 + i * (chip_w + chip_gap)
    chip(slide, kw, cx, kw_y, chip_w, 0.32,
         fill=TEAL, fg=WHITE, size=10, radius=0.5)

footer_brand(slide)
page_number(slide, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Stage 2: Resolve identities (+ ortholog crossover)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 2,
             "Identifying each gene",
             "Turning informal gene names into canonical, unambiguous identifiers.",
             total=TOTAL_STAGES)

# Two parallel transformation rails
rail1_y = 3.05
add_text(slide, "Symbol  →  Canonical fly identity", 0.6, rail1_y, 6.5, 0.4,
         size=14, bold=True, color=INK, font=HEADER_FONT, char_spacing=80)
add_text(slide,
         "Researcher symbols are matched to the official FlyBase identifier "
         "for each gene, plus its primary symbol and known synonyms.",
         0.6, rail1_y + 0.45, 12.1, 0.5,
         size=13, italic=True, color=SLATE_LITE, font=HEADER_FONT)

# Visual chips for rail 1
trans_y = rail1_y + 1.15
chip(slide, "per", 0.6, trans_y, 1.4, 0.55, fill=CREAM_DEEP, fg=INK,
     size=14, radius=0.5)
arrow_glyph(slide, 2.05, trans_y, w=0.55, h=0.55, color=SLATE_LITE, size=22)
chip(slide, "look up canonical identity", 2.6, trans_y, 4.1, 0.55,
     fill=TEAL, fg=WHITE, size=12, radius=0.5)
arrow_glyph(slide, 6.75, trans_y, w=0.55, h=0.55, color=SLATE_LITE, size=22)
# Identity result chip — bigger
add_round_rect(slide, 7.30, trans_y, 5.4, 0.55, CORAL, radius=0.5)
add_text(slide, "period   ·   FBgn0003068   ·   synonyms…",
         7.30, trans_y, 5.4, 0.55,
         size=12, bold=True, color=WHITE, font=BODY_FONT,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# Divider
horizontal_rule(slide, 0.6, 4.85, 12.1, color=CREAM_DEEP, thick=Pt(0.75))

# Rail 2 — optional ortholog crossover
rail2_y = 5.05
add_text(slide,
         "Optional  ·  Cross over to another species",
         0.6, rail2_y, 7.0, 0.4,
         size=14, bold=True, color=INK, font=HEADER_FONT, char_spacing=80)
add_text(slide,
         "If the researcher chooses human or mouse, each fly gene is mapped "
         "to its evolutionary counterpart. One fly gene may map to several "
         "orthologs; unmapped genes are kept and flagged.",
         0.6, rail2_y + 0.45, 12.1, 0.6,
         size=13, italic=True, color=SLATE_LITE, font=HEADER_FONT)

# Branching visual: fly chip → splits into Human / Mouse chips
brail_y = rail2_y + 1.20
chip(slide, "fly gene  (period)", 0.6, brail_y, 3.0, 0.55,
     fill=CREAM_DEEP, fg=INK, size=12, radius=0.5)
arrow_glyph(slide, 3.65, brail_y, w=0.5, h=0.55, color=SLATE_LITE, size=22)
chip(slide, "ortholog mapping", 4.20, brail_y, 2.6, 0.55,
     fill=TEAL, fg=WHITE, size=12, radius=0.5)

# two outputs (human + mouse) right-aligned
chip(slide, "human  ·  PER1, PER2, PER3", 7.10, brail_y - 0.32, 5.6, 0.55,
     fill=CORAL, fg=WHITE, size=12, radius=0.5)
chip(slide, "mouse  ·  Per1, Per2, Per3",  7.10, brail_y + 0.32, 5.6, 0.55,
     fill=CORAL, fg=WHITE, size=12, radius=0.5)

# Two thin connector lines from the mapping chip into the two outputs
def diag_line(x1, y1, x2, y2):
    ln = slide.shapes.add_connector(1,
                                    Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    ln.line.color.rgb = SLATE_LITE
    ln.line.width = Pt(1.25)
    return ln

diag_line(6.82, brail_y + 0.27, 7.10, brail_y - 0.05)
diag_line(6.82, brail_y + 0.27, 7.10, brail_y + 0.59)

footer_brand(slide)
page_number(slide, 4, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Stage 3: Casting the net (multi-source literature)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 3,
             "Casting a wide net",
             "Pulling every paper that might mention each gene, from the right places.",
             total=TOTAL_STAGES)

# Intro line
add_text(slide,
         "The pipeline queries multiple independent databases in parallel. "
         "Different species pull from different curated sources — the right "
         "tool for the right gene.",
         0.6, 2.85, 12.1, 0.7,
         size=13, italic=True, color=SLATE, font=HEADER_FONT)

# Three "lanes" — one per species mode
lanes = [
    ("Fly",   INK,
     ["FlyBase reference tables",
      "PubMed scoped to Drosophila",
      "Europe PMC scoped to Drosophila"]),
    ("Human", TEAL,
     ["HGNC (gene symbols & IDs)",
      "NCBI gene-to-paper links",
      "GeneRIF curated annotations",
      "UniProt reviewed entries",
      "PubMed & Europe PMC scoped to human"]),
    ("Mouse", CORAL,
     ["MGI marker reports",
      "NCBI gene-to-paper links",
      "GeneRIF curated annotations",
      "UniProt reviewed entries",
      "PubMed & Europe PMC scoped to mouse"]),
]

lane_y = 3.7
lane_h = 3.0
lane_w = 4.05
gap = 0.13
for i, (name, color, sources) in enumerate(lanes):
    lx = 0.6 + i * (lane_w + gap)
    # Card
    add_round_rect(slide, lx, lane_y, lane_w, lane_h, WHITE, radius=0.03)
    # Header bar
    add_rect(slide, lx, lane_y, lane_w, 0.55, color)
    add_text(slide, f"{name}  mode",
             lx + 0.25, lane_y, lane_w - 0.5, 0.55,
             size=15, bold=True, color=WHITE, font=HEADER_FONT,
             valign=MSO_ANCHOR.MIDDLE, char_spacing=100)
    # "Sources" label
    add_text(slide, "Sources queried", lx + 0.25, lane_y + 0.7,
             lane_w - 0.5, 0.3,
             size=10, bold=True, color=SLATE_LITE,
             font=BODY_FONT, char_spacing=200)
    # source list
    add_bullets(slide, sources, lx + 0.25, lane_y + 1.05,
                lane_w - 0.5, 1.85, size=11, color=SLATE,
                marker="·", marker_color=color, line_gap=4)

# Bottom band describing what comes out of this stage
out_y = 6.85
add_text(slide,
         "→  Output of this stage:  one combined pool of candidate references for every gene, "
         "with each source tagged.",
         0.6, out_y, 12.1, 0.4,
         size=12, italic=True, bold=True, color=INK, font=HEADER_FONT)

footer_brand(slide)
page_number(slide, 5, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Stage 4: Sorting & filtering
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 4,
             "From a pile of papers to a focused shortlist",
             "Three quick passes turn a noisy pool of references into a tight reading list.",
             total=TOTAL_STAGES)

# Funnel-style three steps
steps = [
    ("Deduplicate",
     "Combine results from every database. The same paper found by three "
     "different sources counts once, but is remembered as a stronger signal.",
     INK_SOFT),
    ("Rank",
     "Papers cited by more sources rise to the top. Newer papers break ties. "
     "The list is then capped at a researcher-chosen size.",
     TEAL),
    ("Filter by keywords",
     "Each paper's title and abstract is scanned for the researcher's keywords. "
     "Papers with no relevant signal are set aside.",
     CORAL),
]

# Visual funnel — three stacked horizontal bars, narrowing
funnel_x = 0.6
funnel_y = 3.0
bar_h = 0.95
gap_y = 0.18
widths = [12.1, 9.5, 7.0]   # narrowing
labels_above = ["Raw pool", "Ranked & capped", "Keyword-relevant"]
counts = ["~30 references / gene", "~15", "~10"]

for i, (w, label, count, (title, desc, color)) in enumerate(zip(widths, labels_above, counts, steps)):
    by = funnel_y + i * (bar_h + gap_y + 0.15)
    cx = funnel_x + (12.1 - w) / 2
    # bar
    add_round_rect(slide, cx, by, w, bar_h, WHITE, radius=0.03)
    # left color marker
    add_rect(slide, cx, by, 0.16, bar_h, color)
    # right count chip — sized to its text, anchored to right edge
    chip_w = 2.1
    chip_x = cx + w - chip_w - 0.18
    chip(slide, count, chip_x, by + 0.27, chip_w, 0.40,
         fill=color, fg=WHITE, size=11, radius=0.5)
    # title and description — width ends well before chip
    text_w = chip_x - (cx + 0.36) - 0.20
    add_text(slide, title, cx + 0.36, by + 0.08, text_w, 0.36,
             size=14, bold=True, color=color, font=HEADER_FONT)
    add_text(slide, desc, cx + 0.36, by + 0.42, text_w, 0.55,
             size=11.5, color=SLATE, font=BODY_FONT)
    # tiny eyebrow above bar
    add_text(slide, label.upper(),
             cx + 0.36, by - 0.25, w - 0.5, 0.22,
             size=9, bold=True, color=SLATE_LITE, font=BODY_FONT,
             char_spacing=200)

footer_brand(slide)
page_number(slide, 6, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Stage 5: Reading & understanding
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 5,
             "Reading the papers",
             "Each surviving reference is read in full, and the gene-relevant evidence is captured in plain language.",
             total=TOTAL_STAGES)

# Two-column layout
# LEFT: How a paper becomes a summary
left_x = 0.6
left_w = 6.0
left_y = 2.95
add_text(slide, "How a paper becomes a summary",
         left_x, left_y, left_w, 0.4,
         size=15, bold=True, color=INK, font=HEADER_FONT)

flow_steps = [
    ("Try to retrieve the full text from any open-access route available.",
     "If full text is not available, fall back to the title and abstract."),
    ("The model reads what it has and writes a focused, evidence-only "
     "summary about this one gene — function, phenotype, and how the paper "
     "supports the researcher's keywords.",
     None),
    ("Summaries that say 'no evidence found' are dropped, so only "
     "informative summaries reach the next stage.",
     None),
]
sy = left_y + 0.6
for i, (main, sub) in enumerate(flow_steps):
    # Step number disc
    add_oval(slide, left_x, sy, 0.42, 0.42, INK)
    add_text(slide, str(i + 1), left_x, sy, 0.42, 0.42,
             size=14, bold=True, color=CREAM, font=HEADER_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, main, left_x + 0.6, sy - 0.02, left_w - 0.6, 0.7,
             size=12, color=SLATE, font=BODY_FONT)
    if sub:
        add_text(slide, sub, left_x + 0.6, sy + 0.65, left_w - 0.6, 0.45,
                 size=11, italic=True, color=SLATE_LITE, font=BODY_FONT)
        sy += 1.30
    else:
        sy += 0.95

# RIGHT: example summary card
right_x = 7.0
right_w = 5.7
right_y = 2.95
add_text(slide, "What one summary looks like",
         right_x, right_y, right_w, 0.4,
         size=15, bold=True, color=INK, font=HEADER_FONT)

# Card
card_y = right_y + 0.55
card_h = 3.7
add_round_rect(slide, right_x, card_y, right_w, card_h, WHITE, radius=0.03)
# header strip
add_rect(slide, right_x, card_y, right_w, 0.42, INK)
add_text(slide, "period  ·  one paper  ·  evidence-only summary",
         right_x + 0.2, card_y, right_w - 0.4, 0.42,
         size=11, bold=True, color=CREAM, font=BODY_FONT,
         valign=MSO_ANCHOR.MIDDLE, char_spacing=80)

# Three labeled blocks: function / phenotype / relevance to keywords
labels = [
    ("FUNCTION",  TEAL,
     "A core component of the fly circadian clock; rhythmically represses "
     "clock-gene transcription on a 24-hour cycle."),
    ("PHENOTYPE", CORAL,
     "Loss-of-function mutants are behaviourally arrhythmic; sleep is "
     "severely disrupted in null animals."),
    ("RELEVANCE", AMBER,
     "Direct, multi-experiment support for both the 'sleep' and 'circadian' "
     "keywords used in this run."),
]
by = card_y + 0.55
for lbl, lc, txt in labels:
    chip(slide, lbl, right_x + 0.2, by, 1.3, 0.32,
         fill=lc, fg=WHITE, size=9, radius=0.5)
    add_text(slide, txt, right_x + 1.6, by - 0.04, right_w - 1.85, 0.85,
             size=11, color=SLATE, font=BODY_FONT)
    by += 1.0

footer_brand(slide)
page_number(slide, 7, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Stage 6: The verdict (classification)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 6,
             "Classifying each gene",
             "All the evidence for one gene is brought together, weighed, and turned into a single verdict.",
             total=TOTAL_STAGES)

# Top row: aggregate + classify visualization
agg_y = 2.95

# Many small "summary" tiles on the left, converging into one classification card
add_text(slide, "Many summaries  →  one decision",
         0.6, agg_y, 6.5, 0.4,
         size=14, bold=True, color=INK, font=HEADER_FONT)

# Six small summary tiles on left
tiles_y = agg_y + 0.55
for i in range(6):
    row = i // 3
    col = i % 3
    tx = 0.6 + col * 1.55
    ty = tiles_y + row * 0.85
    add_round_rect(slide, tx, ty, 1.4, 0.7, WHITE, radius=0.04)
    add_rect(slide, tx, ty, 0.1, 0.7, TEAL)
    add_text(slide, f"summary {i + 1}",
             tx + 0.18, ty + 0.06, 1.2, 0.25,
             size=10, bold=True, color=INK, font=BODY_FONT)
    add_text(slide, "evidence excerpt",
             tx + 0.18, ty + 0.32, 1.2, 0.32,
             size=8.5, italic=True, color=SLATE_LITE, font=BODY_FONT)

# Big arrow
arrow_glyph(slide, 5.5, tiles_y + 0.45, w=1.0, h=0.7, color=CORAL, size=42)

# Verdict card on right
v_x = 6.7
v_w = 6.0
v_y = tiles_y - 0.05
add_round_rect(slide, v_x, v_y, v_w, 1.85, INK, radius=0.04)
add_text(slide, "VERDICT  ·  period",
         v_x + 0.3, v_y + 0.15, v_w - 0.6, 0.35,
         size=11, bold=True, color=AMBER, font=BODY_FONT, char_spacing=200)

# Category line
add_text(slide, "Category",
         v_x + 0.3, v_y + 0.55, 1.4, 0.3,
         size=10, color=INK_MUTED, font=BODY_FONT, char_spacing=120)
add_text(slide, "sleep   ·   circadian",
         v_x + 1.7, v_y + 0.50, v_w - 2.0, 0.4,
         size=18, bold=True, color=CORAL, font=HEADER_FONT)

# Confidence
add_text(slide, "Confidence",
         v_x + 0.3, v_y + 1.00, 1.4, 0.3,
         size=10, color=INK_MUTED, font=BODY_FONT, char_spacing=120)
add_text(slide, "94 / 100",
         v_x + 1.7, v_y + 0.95, 2.5, 0.4,
         size=18, bold=True, color=CREAM, font=HEADER_FONT)

# Rationale
add_text(slide, "Rationale",
         v_x + 0.3, v_y + 1.45, 1.4, 0.3,
         size=10, color=INK_MUTED, font=BODY_FONT, char_spacing=120)
add_text(slide,
         "Multiple papers document arrhythmia and severe sleep disruption.",
         v_x + 1.7, v_y + 1.40, v_w - 2.0, 0.4,
         size=11, italic=True, color=CREAM, font=HEADER_FONT)

# Bottom band: how to read the confidence
horizontal_rule(slide, 0.6, 5.55, 12.1, color=CREAM_DEEP, thick=Pt(0.75))
eyebrow(slide, "How to read the confidence score", 0.6, 5.7, w=6.0)

# Three score bands as colored chips
band_y = 6.05
bands = [
    ("80 – 100",  "Strong, multi-paper support",       TEAL,  WHITE),
    ("50 – 79",   "Partial or indirect evidence",      AMBER, INK),
    ("Below 50",  "Weak or ambiguous — treat as a lead", CORAL, WHITE),
]
band_w = 4.0
for i, (range_, desc, fill, fg) in enumerate(bands):
    bx = 0.6 + i * (band_w + 0.05)
    add_round_rect(slide, bx, band_y, band_w, 0.85, WHITE, radius=0.03)
    # color stripe
    add_rect(slide, bx, band_y, 0.16, 0.85, fill)
    add_text(slide, range_, bx + 0.3, band_y + 0.10, band_w - 0.5, 0.32,
             size=14, bold=True, color=INK, font=HEADER_FONT)
    add_text(slide, desc, bx + 0.3, band_y + 0.45, band_w - 0.5, 0.35,
             size=11, color=SLATE, font=BODY_FONT)

footer_brand(slide)
page_number(slide, 8, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Stage 7: Final report
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
stage_header(slide, 7,
             "Delivering the report",
             "Everything the pipeline learned about your gene set, in one structured Excel workbook.",
             total=TOTAL_STAGES)

# Workbook visual: a stack of three sheet "tabs"
# Left: workbook visual
wb_x = 0.6
wb_y = 3.0
wb_w = 5.6
wb_h = 3.7
add_round_rect(slide, wb_x, wb_y, wb_w, wb_h, WHITE, radius=0.03)

# Tab strip
tab_h = 0.5
tabs = [("Gene Set", INK), ("Classification", TEAL), ("Reference Summaries", CORAL)]
tab_w = wb_w / 3
for i, (name, color) in enumerate(tabs):
    tx = wb_x + i * tab_w
    add_rect(slide, tx, wb_y, tab_w, tab_h, color)
    add_text(slide, name, tx, wb_y, tab_w, tab_h,
             size=11, bold=True, color=WHITE, font=BODY_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
             char_spacing=80)

# Mini "spreadsheet" rows below tabs
row_y = wb_y + tab_h + 0.1
# header row
add_rect(slide, wb_x + 0.15, row_y, wb_w - 0.3, 0.32, CREAM_DEEP)
hdrs = ["Gene", "ID", "Category", "Conf."]
hw = [1.4, 1.4, 1.6, 0.85]
hx = wb_x + 0.25
for h, w in zip(hdrs, hw):
    add_text(slide, h, hx, row_y, w, 0.32,
             size=10, bold=True, color=SLATE, font=BODY_FONT,
             valign=MSO_ANCHOR.MIDDLE)
    hx += w
mock_rows = [
    ("period",     "FBgn0003068", "sleep · circadian", "94", CORAL),
    ("timeless",   "FBgn0014396", "circadian",         "91", TEAL),
    ("redeye",     "FBgn0031722", "sleep",             "78", AMBER),
    ("bruchpilot", "FBgn0025682", "synapse",           "88", TEAL),
    ("sifamide",   "FBgn0039170", "sleep · metab.",    "72", AMBER),
]
for i, (g, fid, cat, conf, col) in enumerate(mock_rows):
    ry = row_y + 0.36 + i * 0.4
    if i % 2 == 0:
        add_rect(slide, wb_x + 0.15, ry, wb_w - 0.3, 0.4, CREAM)
    rx = wb_x + 0.25
    cells = [
        (g,    INK,  True),
        (fid,  SLATE_LITE, False),
        (cat,  col, True),
        (conf, INK, True),
    ]
    for (val, color, bold), w in zip(cells, hw):
        add_text(slide, val, rx, ry, w, 0.4,
                 size=10.5, bold=bold, color=color, font=BODY_FONT,
                 valign=MSO_ANCHOR.MIDDLE)
        rx += w

# RIGHT: what each sheet holds
desc_x = 6.6
desc_w = 6.1
desc_y = 3.0

eyebrow(slide, "What each sheet contains", desc_x, desc_y - 0.05, w=desc_w)

sheet_descs = [
    ("Gene Set", INK,
     "Every gene the pipeline started with. For ortholog runs, each fly gene "
     "appears alongside the species crossover it produced."),
    ("Classification", TEAL,
     "One row per gene with its assigned categories, confidence score, and a "
     "short evidence-backed rationale."),
    ("Reference Summaries", CORAL,
     "Every paper the pipeline read, with title, year, source database, and the "
     "structured evidence summary that fed the classification."),
]
dy = desc_y + 0.4
for name, color, desc in sheet_descs:
    add_rect(slide, desc_x, dy, 0.16, 1.05, color)
    add_text(slide, name, desc_x + 0.32, dy, desc_w - 0.32, 0.32,
             size=14, bold=True, color=color, font=HEADER_FONT)
    add_text(slide, desc, desc_x + 0.32, dy + 0.36, desc_w - 0.32, 0.7,
             size=11.5, color=SLATE, font=BODY_FONT)
    dy += 1.20

footer_brand(slide)
page_number(slide, 9, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Flowchart: Fly mode
# ════════════════════════════════════════════════════════════════════════════
fly_steps = [
    (1, "Read input gene set",
     "Open the input CSV and pull fly gene symbols from the chosen column."),
    (2, "Resolve fly identity",
     "Match each symbol to its FlyBase identifier, primary symbol, "
     "and known synonyms."),
    (3, "Gather references",
     "Query FlyBase, plus PubMed and Europe PMC scoped to Drosophila / "
     "fruit fly."),
    (4, "Merge and deduplicate",
     "Combine results across sources; the same paper found by multiple "
     "sources counts once but ranks higher."),
    (5, "Rank and cap",
     "Sort by source consensus and recency, then keep the top references "
     "up to the run's reference limit."),
    (6, "Filter by keywords",
     "Scan each paper's title and abstract for the researcher's keywords; "
     "drop those with no relevant signal."),
    (7, "Read and summarise",
     "Pull full text where possible. The model writes a focused, "
     "evidence-only summary per paper, per gene."),
    (8, "Classify the gene",
     "Aggregate surviving summaries into one verdict: category, "
     "confidence, and evidence-backed rationale."),
]
organism_flowchart_slide(
    prs, blank,
    slide_num=10, total=TOTAL_SLIDES,
    organism_label="Fly",
    title="How fly mode works",
    subtitle="Fly gene symbols stay in fly throughout — no ortholog crossover.",
    accent=INK,
    steps=fly_steps,
    cols=4,
    input_label="Folder of CSVs with fly gene symbols",
    output_label="<input_name>_classification.xlsx",
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Flowchart: Human mode
# ════════════════════════════════════════════════════════════════════════════
human_steps = [
    (1, "Read input gene set",
     "Open the input CSV and pull fly gene symbols from the chosen column."),
    (2, "Resolve fly identity",
     "Match each symbol to its FlyBase identifier, primary symbol, "
     "and known synonyms."),
    (3, "Map to human orthologs",
     "Cross over to human via DIOPT; one fly gene may map to several "
     "human orthologs, all retained."),
    (4, "Build human gene set",
     "Resolve each human ortholog to its canonical symbol, IDs, and "
     "synonyms. Unmapped fly genes are kept and flagged."),
    (5, "Gather references",
     "Query NCBI gene-to-paper, GeneRIF, UniProt, plus PubMed and "
     "Europe PMC scoped to Homo sapiens."),
    (6, "Merge and deduplicate",
     "Combine results across sources; the same paper found by multiple "
     "sources counts once but ranks higher."),
    (7, "Rank and cap",
     "Sort by source consensus and recency, then keep the top references "
     "up to the run's reference limit."),
    (8, "Filter by keywords",
     "Scan each paper's title and abstract for the researcher's keywords; "
     "drop those with no relevant signal."),
    (9, "Read and summarise",
     "Pull full text where possible. The model writes a focused, "
     "evidence-only summary per paper, per gene."),
    (10, "Classify the gene",
     "Aggregate surviving summaries into one verdict: category, "
     "confidence, and evidence-backed rationale."),
]
organism_flowchart_slide(
    prs, blank,
    slide_num=11, total=TOTAL_SLIDES,
    organism_label="Human",
    title="How human mode works",
    subtitle="Fly genes are mapped to their human orthologs before any literature is collected.",
    accent=TEAL,
    steps=human_steps,
    cols=5,
    input_label="Folder of CSVs with fly gene symbols",
    output_label="<input_name>_human_classification.xlsx  +  <input_name>_human.csv",
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Flowchart: Mouse mode
# ════════════════════════════════════════════════════════════════════════════
mouse_steps = [
    (1, "Read input gene set",
     "Open the input CSV and pull fly gene symbols from the chosen column."),
    (2, "Resolve fly identity",
     "Match each symbol to its FlyBase identifier, primary symbol, "
     "and known synonyms."),
    (3, "Map to mouse orthologs",
     "Cross over to mouse via DIOPT; one fly gene may map to several "
     "mouse orthologs, all retained."),
    (4, "Build mouse gene set",
     "Resolve each mouse ortholog to its canonical symbol, IDs, and "
     "synonyms via MGI. Unmapped fly genes are kept and flagged."),
    (5, "Gather references",
     "Query NCBI gene-to-paper, GeneRIF, UniProt, plus PubMed and "
     "Europe PMC scoped to Mus musculus."),
    (6, "Merge and deduplicate",
     "Combine results across sources; the same paper found by multiple "
     "sources counts once but ranks higher."),
    (7, "Rank and cap",
     "Sort by source consensus and recency, then keep the top references "
     "up to the run's reference limit."),
    (8, "Filter by keywords",
     "Scan each paper's title and abstract for the researcher's keywords; "
     "drop those with no relevant signal."),
    (9, "Read and summarise",
     "Pull full text where possible. The model writes a focused, "
     "evidence-only summary per paper, per gene."),
    (10, "Classify the gene",
     "Aggregate surviving summaries into one verdict: category, "
     "confidence, and evidence-backed rationale."),
]
organism_flowchart_slide(
    prs, blank,
    slide_num=12, total=TOTAL_SLIDES,
    organism_label="Mouse",
    title="How mouse mode works",
    subtitle="Fly genes are mapped to their mouse orthologs before any literature is collected.",
    accent=CORAL,
    steps=mouse_steps,
    cols=5,
    input_label="Folder of CSVs with fly gene symbols",
    output_label="<input_name>_mouse_classification.xlsx  +  <input_name>_mouse.csv",
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Real example: Tx-Omics_Conserved input and ortholog expansion
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
example_slide_header(
    slide,
    "Ortholog walkthrough  ·  Real run",
    "Example run: Tx-Omics_Conserved",
    "The input folder contains one fly gene set and completed human + mouse ortholog outputs.",
    slide_num=13,
    total=TOTAL_SLIDES,
    accent=TEAL_DEEP,
)

stat_card(slide, 0.6, 2.85, 2.55, 1.55,
          value="8", label="fly genes in input",
          note="from fl.AI-Dmel-gene-hits.csv", accent=INK)
stat_card(slide, 3.35, 2.85, 2.55, 1.55,
          value="16", label="human ortholog genes",
          note="plus 2 unmapped fly rows", accent=TEAL)
stat_card(slide, 6.10, 2.85, 2.55, 1.55,
          value="16", label="mouse ortholog genes",
          note="plus 2 unmapped fly rows", accent=CORAL)

# Input mini spreadsheet
input_x = 0.6
input_y = 4.75
add_round_rect(slide, input_x, input_y, 4.2, 1.85, WHITE, radius=0.03)
add_rect(slide, input_x, input_y, 4.2, 0.38, INK)
add_text(slide, "INPUT CSV", input_x + 0.16, input_y + 0.07, 1.3, 0.25,
         size=9, bold=True, color=AMBER, font=BODY_FONT, char_spacing=180)
add_text(slide, "ext_gene", input_x + 0.20, input_y + 0.48, 1.35, 0.25,
         size=9.5, bold=True, color=SLATE, font=BODY_FONT)
add_text(slide, "flybase_gene_id", input_x + 1.7, input_y + 0.48, 1.8, 0.25,
         size=9.5, bold=True, color=SLATE, font=BODY_FONT)
for i, (gene, fid) in enumerate(EXAMPLE_INPUT_GENES):
    row_y = input_y + 0.78 + i * 0.12
    add_text(slide, gene, input_x + 0.20, row_y, 1.25, 0.14,
             size=8.2, bold=True, color=INK, font=BODY_FONT)
    add_text(slide, fid, input_x + 1.70, row_y, 1.55, 0.14,
             size=8.2, color=SLATE, font=BODY_FONT)

# Ortholog file fan-out
fan_x = 5.15
fan_y = 4.75
add_text(slide, "Ortholog run settings", fan_x, fan_y - 0.10, 7.55, 0.30,
         size=12.5, bold=True, color=INK, font=HEADER_FONT)
for arrow_y, arrow_color in [(fan_y + 0.64, TEAL), (fan_y + 1.42, CORAL)]:
    line = slide.shapes.add_connector(
        1, Inches(input_x + 4.35), Inches(input_y + 0.95),
        Inches(fan_x - 0.15), Inches(arrow_y))
    line.line.color.rgb = arrow_color
    line.line.width = Pt(1.25)
    add_text(slide, "→", fan_x - 0.42, arrow_y - 0.18, 0.24, 0.30,
             size=14, bold=True, color=arrow_color, font=BODY_FONT,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
for i, (label, text, accent) in enumerate([
    ("Human output", "fl.AI-Dmel-gene-hits_human.csv\nfl.AI-Dmel-gene-hits_human_classification.xlsx", TEAL),
    ("Mouse output", "fl.AI-Dmel-gene-hits_mouse.csv\nfl.AI-Dmel-gene-hits_mouse_classification.xlsx", CORAL),
]):
    y = fan_y + 0.35 + i * 0.78
    add_round_rect(slide, fan_x, y, 7.55, 0.58, WHITE, radius=0.04)
    add_rect(slide, fan_x, y, 0.14, 0.58, accent)
    add_text(slide, label, fan_x + 0.32, y + 0.08, 1.8, 0.20,
             size=10.5, bold=True, color=accent, font=BODY_FONT)
    add_text(slide, text, fan_x + 2.10, y + 0.07, 5.15, 0.45,
             size=9.5, color=SLATE, font=BODY_FONT)

add_round_rect(slide, fan_x, 6.45, 7.55, 0.45, INK, radius=0.04)
add_text(slide, "KEYWORDS", fan_x + 0.25, 6.55, 1.2, 0.24,
         size=9, bold=True, color=AMBER, font=BODY_FONT, char_spacing=180)
add_text(slide, "sleep, circadian  ·  reference limit 50  ·  input gene column: ext_gene",
         fan_x + 1.55, 6.54, 5.65, 0.26,
         size=10.5, italic=True, color=CREAM, font=HEADER_FONT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Human ortholog results
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
example_slide_header(
    slide,
    "Ortholog walkthrough  ·  Human result",
    "Fly genes expanded into human literature targets",
    "Human mode converts the 8 fly genes into 16 mapped human orthologs, then classifies the evidence for sleep and circadian relevance.",
    slide_num=14,
    total=TOTAL_SLIDES,
    accent=TEAL,
)

stat_card(slide, 0.6, 2.82, 2.25, 1.35,
          value="16", label="mapped genes", note="from 6 mapped fly genes", accent=TEAL)
stat_card(slide, 3.05, 2.82, 2.25, 1.35,
          value="7", label="classified genes", note="assigned to at least one bucket", accent=CORAL)
stat_card(slide, 5.50, 2.82, 2.25, 1.35,
          value="9", label="unclassified", note="reviewed but no matching bucket", accent=SLATE_LITE)

add_round_rect(slide, 0.6, 4.45, 3.45, 1.55, INK, radius=0.04)
add_text(slide, "HUMAN BUCKETS", 0.85, 4.62, 2.9, 0.28,
         size=10, bold=True, color=AMBER, font=BODY_FONT, char_spacing=200)
add_text(slide, "4", 0.85, 5.00, 0.55, 0.38,
         size=23, bold=True, color=CORAL, font=HEADER_FONT)
add_text(slide, "sleep + circadian", 1.45, 5.06, 2.3, 0.28,
         size=12.5, bold=True, color=CREAM, font=HEADER_FONT)
add_text(slide, "3", 0.85, 5.43, 0.55, 0.38,
         size=23, bold=True, color=TEAL, font=HEADER_FONT)
add_text(slide, "circadian only", 1.45, 5.49, 2.3, 0.28,
         size=12.5, bold=True, color=CREAM, font=HEADER_FONT)

add_round_rect(slide, 4.35, 4.45, 3.40, 1.55, WHITE, radius=0.04)
add_text(slide, "Strongest pattern", 4.60, 4.62, 2.9, 0.28,
         size=12, bold=True, color=TEAL, font=HEADER_FONT)
add_text(slide,
         "The na / unc79 / unc80 fly genes converge on the human NALCN channelosome, producing the highest-confidence sleep + circadian calls.",
         4.60, 4.98, 2.85, 0.78,
         size=10.5, color=SLATE, font=BODY_FONT)

add_simple_table(
    slide,
    x=8.05, y=2.82,
    col_ws=[1.15, 1.25, 1.70, 0.55],
    row_h=0.46,
    headers=["Human", "Fly source", "Bucket", "Conf."],
    rows=HUMAN_RESULT_ROWS,
    accent=TEAL,
    font_size=9.0,
    header_size=9.1,
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Mouse ortholog results
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
example_slide_header(
    slide,
    "Ortholog walkthrough  ·  Mouse result",
    "The same fly set produces a richer mouse result set",
    "Mouse mode uses the same input and keywords, but the mouse literature supports more classified sleep/circadian calls.",
    slide_num=15,
    total=TOTAL_SLIDES,
    accent=CORAL,
)

stat_card(slide, 0.6, 2.82, 2.25, 1.35,
          value="16", label="mapped genes", note="from 6 mapped fly genes", accent=CORAL)
stat_card(slide, 3.05, 2.82, 2.25, 1.35,
          value="11", label="classified genes", note="assigned to at least one bucket", accent=TEAL)
stat_card(slide, 5.50, 2.82, 2.25, 1.35,
          value="5", label="unclassified", note="reviewed but no matching bucket", accent=SLATE_LITE)

add_round_rect(slide, 0.6, 4.45, 3.45, 1.70, INK, radius=0.04)
add_text(slide, "MOUSE BUCKETS", 0.85, 4.62, 2.9, 0.28,
         size=10, bold=True, color=AMBER, font=BODY_FONT, char_spacing=200)
for i, (num, label, accent) in enumerate([
    ("5", "circadian only", TEAL),
    ("4", "sleep + circadian", CORAL),
    ("2", "sleep only", AMBER),
]):
    y = 4.95 + i * 0.36
    add_text(slide, num, 0.85, y, 0.50, 0.32,
             size=18, bold=True, color=accent, font=HEADER_FONT)
    add_text(slide, label, 1.42, y + 0.04, 2.2, 0.25,
             size=11.2, bold=True, color=CREAM, font=HEADER_FONT)

add_round_rect(slide, 4.35, 4.45, 3.40, 1.70, WHITE, radius=0.04)
add_text(slide, "Mouse adds signal", 4.60, 4.62, 2.9, 0.28,
         size=12, bold=True, color=CORAL, font=HEADER_FONT)
add_text(slide,
         "Nalcn and Tph2 become the strongest calls. Mouse-specific evidence also pulls in Galr1 as sleep-only and several rumpel orthologs.",
         4.60, 4.98, 2.85, 0.88,
         size=10.5, color=SLATE, font=BODY_FONT)

add_simple_table(
    slide,
    x=8.05, y=2.82,
    col_ws=[1.15, 1.25, 1.70, 0.55],
    row_h=0.42,
    headers=["Mouse", "Fly source", "Bucket", "Conf."],
    rows=MOUSE_RESULT_ROWS,
    accent=CORAL,
    font_size=8.7,
    header_size=8.8,
)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Human vs mouse takeaways
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, INK)
add_rect(slide, 0, 0, 13.33, 0.16, TEAL)
add_rect(slide, 0, 7.34, 13.33, 0.16, CORAL)

eyebrow(slide, "Ortholog walkthrough  ·  Interpretation", 0.6, 0.65, w=8.5,
        color=AMBER)
add_text(slide, "Same fly genes, different evidence landscapes",
         0.6, 1.05, 12.1, 0.72,
         size=35, bold=True, color=CREAM, font=HEADER_FONT)
add_text(slide,
         "Ortholog mode is not just name translation: the species-specific literature changes which targets get strong evidence.",
         0.6, 1.85, 12.1, 0.42,
         size=14.5, italic=True, color=INK_MUTED, font=HEADER_FONT)

comparison_rows = [
    ("na (fly gene)", "NALCN  ·  92  ·  sleep + circadian", "Nalcn  ·  96  ·  sleep + circadian"),
    ("Trhn", "TH  ·  63  ·  sleep + circadian", "Tph2  ·  95  ·  sleep + circadian"),
    ("AstA-R2", "GALR1/2/3  ·  circadian", "Kiss1r  ·  91  ·  circadian;  Galr1  ·  90  ·  sleep"),
    ("BomBc2, SIFa", "kept as unmapped rows", "kept as unmapped rows"),
]

x0 = 0.75
y0 = 2.75
col_w = [1.55, 5.05, 5.05]
headers = [("Fly source", AMBER), ("Human result", TEAL), ("Mouse result", CORAL)]
cx = x0
for (h, color), w in zip(headers, col_w):
    add_rect(slide, cx, y0 + 0.28, w - 0.25, 0.05, color)
    add_text(slide, h, cx, y0, w - 0.15, 0.28,
             size=12.2, bold=True, color=CREAM, font=BODY_FONT,
             char_spacing=150)
    cx += w

for i, (fly, human, mouse) in enumerate(comparison_rows):
    row_y = y0 + 0.52 + i * 0.68
    add_round_rect(slide, x0, row_y, sum(col_w) + 0.25, 0.52,
                   INK_SOFT if i % 2 == 0 else RGBColor(0x15, 0x2A, 0x43),
                   line=RGBColor(0x45, 0x56, 0x6F), radius=0.03)
    cx = x0 + 0.18
    add_text(slide, fly, cx, row_y + 0.12, col_w[0] - 0.28, 0.24,
             size=11.8, bold=True, color=CREAM, font=BODY_FONT)
    cx += col_w[0]
    add_text(slide, human, cx, row_y + 0.12, col_w[1] - 0.25, 0.24,
             size=11.6, bold=True, color=WHITE, font=BODY_FONT)
    cx += col_w[1]
    add_text(slide, mouse, cx, row_y + 0.12, col_w[2] - 0.25, 0.24,
             size=11.6, bold=True, color=WHITE, font=BODY_FONT)

add_round_rect(slide, 0.75, 6.10, 11.8, 0.80, CREAM, radius=0.04)
add_text(slide, "Takeaway", 1.05, 6.24, 1.2, 0.25,
         size=12, bold=True, color=CORAL, font=BODY_FONT, char_spacing=150)
add_text(slide,
         "The workbook preserves both views: mapping rows explain how fly genes crossed species, while classification rows show which orthologs earned sleep or circadian evidence.",
         2.05, 6.20, 10.05, 0.34,
         size=12.3, italic=True, color=INK, font=HEADER_FONT)

add_text(slide, "fl.AI  ·  Tx-Omics_Conserved example  ·  human + mouse ortholog runs",
         0.6, 7.05, 8.5, 0.22,
         size=9.5, color=INK_MUTED, font=BODY_FONT, char_spacing=150)
page_number(slide, 16, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — The whole journey, one example
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.33, 7.5, INK)
# small accent bands top/bottom
add_rect(slide, 0, 0, 13.33, 0.16, CORAL)
add_rect(slide, 0, 7.34, 13.33, 0.16, TEAL)

eyebrow(slide, "The whole journey, in one example", 0.6, 0.65, w=8.0,
        color=AMBER)
add_text(slide, "One gene.  Seven stages.",
         0.6, 1.05, 12.1, 0.85,
         size=42, bold=True, color=CREAM, font=HEADER_FONT)
add_text(slide,
         "Following a single gene from researcher symbol to final verdict.",
         0.6, 1.95, 12.1, 0.4,
         size=15, italic=True, color=INK_MUTED, font=HEADER_FONT)

# Long horizontal flow with seven stops
flow_y = 3.3
flow_x = 0.6
flow_w = 12.1
n = 7
slot_w = flow_w / n

# Track line
trk_y = flow_y + 0.85
ln = slide.shapes.add_connector(1,
                                Inches(flow_x + slot_w / 2),
                                Inches(trk_y),
                                Inches(flow_x + flow_w - slot_w / 2),
                                Inches(trk_y))
ln.line.color.rgb = SLATE_LITE
ln.line.width = Pt(1.5)

stops = [
    ("Input",      "per",                   CREAM_DEEP, INK),
    ("Resolve",    "FBgn0003068",           TEAL,       WHITE),
    ("Gather",     "~30 candidate papers",  TEAL,       WHITE),
    ("Filter",     "10 keyword-relevant",   TEAL,       WHITE),
    ("Read",       "6 evidence summaries",  TEAL,       WHITE),
    ("Classify",   "sleep · circadian",     CORAL,      WHITE),
    ("Deliver",    "row in the report",     CORAL,      WHITE),
]

# Dots on the track
for i, (stop, payload, fill, fg) in enumerate(stops):
    cx = flow_x + i * slot_w + slot_w / 2
    # dot
    add_oval(slide, cx - 0.10, trk_y - 0.10, 0.20, 0.20, fill)
    # eyebrow above
    add_text(slide, stop.upper(), cx - 1.0, flow_y, 2.0, 0.3,
             size=10, bold=True, color=AMBER, font=BODY_FONT,
             align=PP_ALIGN.CENTER, char_spacing=200)
    # number
    add_text(slide, f"0{i + 1}", cx - 1.0, flow_y + 0.30, 2.0, 0.45,
             size=20, bold=True, color=CREAM, font=HEADER_FONT,
             align=PP_ALIGN.CENTER)
    # payload (below the dot)
    add_text(slide, payload, cx - 1.0, trk_y + 0.30, 2.0, 0.65,
             size=11, italic=True, color=INK_MUTED, font=HEADER_FONT,
             align=PP_ALIGN.CENTER)

# Bottom — summary band
band_y = 5.5
add_round_rect(slide, 0.6, band_y, 12.1, 1.4, INK_SOFT, radius=0.04)
add_text(slide, "WHAT THE RESEARCHER GETS BACK",
         0.85, band_y + 0.15, 11.6, 0.3,
         size=11, bold=True, color=AMBER, font=BODY_FONT, char_spacing=200)

cols = [
    ("A complete record",       "Every paper the pipeline considered, kept or dropped, with the source it came from."),
    ("A reasoned verdict",      "A category, a confidence score, and a short rationale per gene."),
    ("A reusable starting point", "Re-running with new keywords reuses the work already done — no duplicate searching."),
]
col_w = (12.1 - 0.5) / 3
for i, (h, body) in enumerate(cols):
    cx = 0.85 + i * col_w
    add_text(slide, h, cx, band_y + 0.50, col_w - 0.2, 0.32,
             size=13, bold=True, color=CREAM, font=HEADER_FONT)
    add_text(slide, body, cx, band_y + 0.85, col_w - 0.2, 0.5,
             size=11, color=INK_MUTED, font=BODY_FONT)

# Footer
add_text(slide, "fl.AI  ·  Allada Lab  ·  University of Michigan  ·  2026",
         0.6, 7.05, 12.1, 0.22,
         size=10, color=INK_MUTED, font=BODY_FONT,
         align=PP_ALIGN.CENTER, char_spacing=200)
page_number(slide, 17, TOTAL_SLIDES)

# ── Save ─────────────────────────────────────────────────────────────────────
import os
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DECK_DIR   = os.path.normpath(os.path.join(SCRIPT_DIR, "..", "Decks"))
os.makedirs(DECK_DIR, exist_ok=True)
out_path = os.path.join(DECK_DIR, "fl.AI_Data_Flow_Deck.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
